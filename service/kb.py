import uuid
import math
import io
import json
import zipfile
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path
import re
from collections import Counter

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_pdf_text

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_pdf_text

from dao.kb_dao import (
    create_kb,
    update_kb,
    delete_kb,
    list_kb,
    get_kb,
    create_doc,
    update_doc,
    delete_doc,
    list_docs,
    get_doc,
    upsert_doc_embeddings,
    list_doc_embeddings,
    search_doc_embeddings_by_vector,
    search_docs_fulltext,
)
from models.kb import (
    KnowledgeBase,
    KnowledgeBaseCreate,
    KnowledgeBaseUpdate,
    KnowledgeDocument,
    KnowledgeDocumentCreate,
    KnowledgeDocumentUpdate,
    KnowledgeQAReply,
)
from service.openai_service import chat_completion, create_embeddings


def _now_ms() -> int:
    return int(datetime.utcnow().timestamp() * 1000)


def _cosine_similarity(a: List[float], b: List[float]) -> float:
    if len(a) != len(b):
        return 0.0
    dot = sum(x * y for x, y in zip(a, b))
    norm_a = math.sqrt(sum(x * x for x in a))
    norm_b = math.sqrt(sum(y * y for y in b))
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return dot / (norm_a * norm_b)


def _get_owned_kb(kb_uuid: str, owner_uuid: str) -> Optional[KnowledgeBase]:
    kb_data = get_kb(kb_uuid, owner_uuid=owner_uuid)
    if not kb_data:
        return None
    return KnowledgeBase(**kb_data)


def get_owned_kb(kb_uuid: str, owner_uuid: str) -> Optional[KnowledgeBase]:
    return _get_owned_kb(kb_uuid, owner_uuid)


# ==== kb ====


def create_kb_service(owner_uuid: str, req: KnowledgeBaseCreate) -> KnowledgeBase:
    kb = KnowledgeBase(
        uuid=str(uuid.uuid4()),
        name=req.name,
        description=req.description,
        owner_uuid=owner_uuid,
        create_at=_now_ms(),
        update_at=_now_ms(),
    )
    create_kb(kb.dict())
    return kb


def update_kb_service(owner_uuid: str, uuid_: str, req: KnowledgeBaseUpdate) -> Optional[KnowledgeBase]:
    kb = _get_owned_kb(uuid_, owner_uuid)
    if not kb:
        return None

    fields: Dict[str, Any] = {}
    if req.name is not None:
        fields["name"] = req.name
    if req.description is not None:
        fields["description"] = req.description
    if not fields:
        return kb

    fields["update_at"] = _now_ms()
    update_kb(uuid_, fields, owner_uuid=owner_uuid)
    kb_dict = kb.dict()
    kb_dict.update(fields)
    return KnowledgeBase(**kb_dict)


def delete_kb_service(owner_uuid: str, uuid_: str) -> bool:
    kb = _get_owned_kb(uuid_, owner_uuid)
    if not kb:
        return False
    delete_kb(uuid_)
    return True


def list_kb_service(owner_uuid: str, page: int, size: int) -> Dict[str, Any]:
    return list_kb(page, size, owner_uuid)


# ==== doc ====


def create_doc_service(
    owner_uuid: str, kb_uuid: str, req: KnowledgeDocumentCreate
) -> Optional[KnowledgeDocument]:
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return None

    doc = KnowledgeDocument(
        uuid=str(uuid.uuid4()),
        kb_uuid=kb_uuid,
        title=req.title,
        content=req.content,
        create_at=_now_ms(),
        update_at=_now_ms(),
    )
    create_doc(doc.dict())

    # generate embedding and write into
    _generate_and_store_embeddings_for_doc(doc)

    return doc


def update_doc_service(
    owner_uuid: str, uuid_: str, req: KnowledgeDocumentUpdate
) -> Optional[KnowledgeDocument]:
    doc_data = get_doc(uuid_)
    if not doc_data:
        return None
    kb_uuid = doc_data.get("kb_uuid")
    if not kb_uuid or not _get_owned_kb(kb_uuid, owner_uuid):
        return None

    fields: Dict[str, Any] = {}
    if req.title is not None:
        fields["title"] = req.title
    if req.content is not None:
        fields["content"] = req.content
    if not fields:
        return KnowledgeDocument(**doc_data)

    fields["update_at"] = _now_ms()
    update_doc(uuid_, fields)
    doc_data.update(fields)
    doc = KnowledgeDocument(**doc_data)

    # if content has changed, regenerate embedding
    if req.content is not None:
        _generate_and_store_embeddings_for_doc(doc)

    return doc


def delete_doc_service(owner_uuid: str, uuid_: str) -> bool:
    doc_data = get_doc(uuid_)
    if not doc_data or not _get_owned_kb(doc_data.get("kb_uuid", ""), owner_uuid):
        return False
    delete_doc(uuid_)
    return True


def list_docs_service(owner_uuid: str, kb_uuid: str, page: int, size: int) -> Dict[str, Any]:
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return {"total": 0, "list": []}
    return list_docs(kb_uuid, page, size)


def _chunk_text(content: str, max_chars: int = 400) -> List[str]:
    """
    currently kb is mainly used for ES query, not fed to OpenAI.
    keep a simple chunking function for future extension.
    """
    content = content.strip()
    if not content:
        return []
    chunks: List[str] = []
    start = 0
    while start < len(content):
        end = min(start + max_chars, len(content))
        chunks.append(content[start:end])
        start = end
    return chunks


def _generate_and_store_embeddings_for_doc(doc: KnowledgeDocument) -> None:
    chunks = _chunk_text(doc.content)
    if not chunks:
        return

    vectors: List[Dict[str, Any]] = []
    for chunk in chunks:
        embedding = create_embeddings(chunk)
        vectors.append(
            {
                "uuid": str(uuid.uuid4()),
                "chunk": chunk,
                "embedding": embedding,
                "create_at": _now_ms(),
            }
        )
    upsert_doc_embeddings(doc.kb_uuid, doc.uuid, vectors)


def qa_service(owner_uuid: str, kb_uuid: str, question: str, top_k: int = 3) -> Optional[KnowledgeQAReply]:
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return None

    context_chunks = _retrieve_context_chunks(kb_uuid, question, top_k)
    messages = _build_messages_with_context(question, context_chunks)
    answer = chat_completion(messages)

    # write current Q&A into kb, and generate vector for the answer
    save_qa_to_kb(kb_uuid, question, answer)

    context_texts = [item["chunk"] for item in context_chunks]
    return KnowledgeQAReply(answer=answer, context=context_texts)


def save_qa_to_kb(kb_uuid: str, question: str, answer: str) -> None:
    """
    write current Q&A into kb, and generate vector for the answer
    """
    doc = KnowledgeDocument(
        uuid=str(uuid.uuid4()),
        kb_uuid=kb_uuid,
        title=question[:50],
        content=f"Q: {question}\n\nA: {answer}",
        create_at=_now_ms(),
        update_at=_now_ms(),
    )
    create_doc(doc.dict())

    # only generate embedding for the answer text
    embedding = create_embeddings(answer)
    upsert_doc_embeddings(
        kb_uuid,
        doc.uuid,
        [
            {
                "uuid": str(uuid.uuid4()),
                "chunk": answer,
                "embedding": embedding,
                "create_at": _now_ms(),
            }
        ],
    )


def semantic_search_service(owner_uuid: str, kb_uuid: str, query: str, top_k: int = 5) -> Optional[List[Dict[str, Any]]]:
    """
    do vector semantic search for the specified kb:
    - generate embedding for the query
    - fetch all vectors under the kb from kb_doc_embed_index
    - calculate cosine similarity, return top_k chunks + scores
    """
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return None

    query_vector = create_embeddings(query)
    results: List[Dict[str, Any]] = []
    try:
        results = search_doc_embeddings_by_vector(kb_uuid, query_vector, top_k)
    except Exception as exc:  # pylint: disable=broad-except
        print(f"[WARN] ES vector search failed, falling back to local scoring: {exc}")
        vectors = list_doc_embeddings(kb_uuid)
        results = _score_vectors_locally(
            vectors,
            query_vector,
            top_k=top_k,
            score_threshold=0.0,
        )

    formatted: List[Dict[str, Any]] = []
    for item in results[:top_k]:
        formatted.append(
            {
                "kb_uuid": item.get("kb_uuid"),
                "doc_uuid": item.get("doc_uuid"),
                "chunk": item.get("chunk", ""),
                "score": item.get("score", 0.0),
            }
        )
    return formatted


def fulltext_search_service(
    owner_uuid: str,
    kb_uuid: str,
    query: str,
    top_k: int = 5,
) -> Optional[List[Dict[str, Any]]]:
    """
    Keyword-based full-text search with ES highlighting.
    """
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return None
    return search_docs_fulltext(kb_uuid, query, top_k)


def import_kb_file_service(
    owner_uuid: str,
    kb_uuid: str,
    filename: str,
    file_bytes: bytes,
) -> Optional[Dict[str, Any]]:
    """
    Parse uploaded file(s) and insert docs into KB.
    Supports markdown/txt/csv/docx/pptx/pdf.
    """
    if not _get_owned_kb(kb_uuid, owner_uuid):
        return None

    docs = _extract_docs_from_upload(filename, file_bytes)
    summary = {
        "total": len(docs),
        "success": 0,
        "failed": 0,
        "errors": [],
    }

    for idx, payload in enumerate(docs, start=1):
        title = (payload.get("title") or f"Imported {idx}").strip()
        content = (payload.get("content") or "").strip()
        if not content:
            summary["failed"] += 1
            if len(summary["errors"]) < 20:
                summary["errors"].append(f"{title or 'Document'} has empty content, skipped")
            continue
        try:
            create_doc_service(
                owner_uuid,
                kb_uuid,
                KnowledgeDocumentCreate(title=title or f"Imported {idx}", content=content),
            )
            summary["success"] += 1
        except Exception as exc:  # pylint: disable=broad-except
            summary["failed"] += 1
            if len(summary["errors"]) < 20:
                summary["errors"].append(f"{title[:50] or 'Document'}: {exc}")

    return summary


def _retrieve_context_chunks(
    kb_uuid: str,
    question: str,
    top_k: int = 3,
    score_threshold: float = 0.2,
) -> List[Dict[str, Any]]:
    """
    Retrieve top_k most relevant chunks from KB embeddings.
    Falls back gracefully if no embeddings exist or ES vector search fails.
    """
    query_vector = create_embeddings(question)
    scored: List[Dict[str, Any]] = []

    try:
        scored = [
            item
            for item in search_doc_embeddings_by_vector(
                kb_uuid,
                query_vector,
                top_k=max(top_k, 5),
            )
            if item.get("score", 0.0) >= score_threshold
        ]
    except Exception as exc:  # pylint: disable=broad-except
        print(f"[WARN] ES vector search failed, fallback to local scoring: {exc}")
        vectors = list_doc_embeddings(kb_uuid)
        scored = _score_vectors_locally(
            vectors,
            query_vector,
            top_k=max(top_k, 5),
            score_threshold=score_threshold,
        )

    return scored[:top_k]


def _build_messages_with_context(
    question: str, context_chunks: List[Dict[str, Any]]
) -> List[Dict[str, str]]:
    """
    Construct system/user messages. Encourage the model to use KB context when available.
    """
    base_instruction = (
        "You are a helpful assistant. Use the provided knowledge base snippets when they are relevant. "
        "If the context does not contain sufficient information, clearly state that you are not sure "
        "instead of fabricating details."
    )

    messages: List[Dict[str, str]] = [{"role": "system", "content": base_instruction}]

    if context_chunks:
        context_text = "\n\n".join(
            f"[Score {item['score']:.2f}] {item['chunk']}" for item in context_chunks
        )
        messages.append({"role": "system", "content": f"Knowledge base context:\n{context_text}"})

    messages.append({"role": "user", "content": question})
    return messages


def _score_vectors_locally(
    vectors: List[Dict[str, Any]],
    query_vector: List[float],
    top_k: int,
    score_threshold: float,
) -> List[Dict[str, Any]]:
    """Fallback cosine scoring executed in Python."""
    scored: List[Dict[str, Any]] = []
    for item in vectors:
        emb = item.get("embedding") or []
        if not emb:
            continue
        score = _cosine_similarity(query_vector, emb)
        if score < score_threshold:
            continue
        scored.append(
            {
                "kb_uuid": item.get("kb_uuid"),
                "doc_uuid": item.get("doc_uuid"),
                "chunk": item.get("chunk", ""),
                "score": score,
            }
        )
    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:top_k]


def export_kb_service(owner_uuid: str, kb_uuid: str) -> Optional[Dict[str, Any]]:
    """
    Bundle kb metadata, documents, and embeddings into a zip for download.
    """
    kb = _get_owned_kb(kb_uuid, owner_uuid)
    if not kb:
        return None

    kb_data = kb.dict()
    docs = _fetch_all_docs(kb_uuid)
    embeddings = list_doc_embeddings(kb_uuid)

    bundle = {
        "kb": kb_data,
        "documents": docs,
        "embeddings": embeddings,
        "exported_at": _now_ms(),
    }

    memory_file = io.BytesIO()
    with zipfile.ZipFile(memory_file, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("kb.json", json.dumps(kb_data, ensure_ascii=False, indent=2))
        zf.writestr("docs.json", json.dumps(docs, ensure_ascii=False))
        zf.writestr(
            "embeddings.json",
            json.dumps(embeddings, ensure_ascii=False),
        )
        zf.writestr(
            "bundle.json",
            json.dumps(bundle, ensure_ascii=False),
        )

    memory_file.seek(0)
    safe_name = re.sub(r"[^a-zA-Z0-9_-]", "-", kb_data.get("name", "kb"))
    filename = f"{safe_name or 'kb'}-{kb_uuid[:8]}.zip"
    return {"filename": filename, "content": memory_file.read()}


def _fetch_all_docs(kb_uuid: str, page_size: int = 200) -> List[Dict[str, Any]]:
    docs: List[Dict[str, Any]] = []
    page = 1
    total = 0
    while True:
        batch = list_docs(kb_uuid, page, page_size)
        items = batch.get("list", [])
        total = batch.get("total", 0)
        docs.extend(items)
        if len(docs) >= total or not items:
            break
        page += 1
    return docs


def _extract_docs_from_upload(filename: str, payload: bytes) -> List[Dict[str, str]]:
    suffix = Path((filename or "")).suffix.lower()
    if suffix in {".md", ".markdown"}:
        return _parse_markdown_documents(_decode_text(payload))
    if suffix in {".txt", ""}:
        return _parse_plain_text(_decode_text(payload), filename)
    if suffix == ".csv":
        return _parse_csv_documents(payload, filename)
    if suffix == ".docx":
        return _parse_docx_documents(payload, filename)
    if suffix == ".pptx":
        return _parse_pptx_documents(payload, filename)
    if suffix == ".pdf":
        return _parse_pdf_document(payload, filename)
    raise ValueError("Unsupported file format. Use markdown/txt, csv, docx, pptx, or pdf.")


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("latin-1", errors="ignore")


def _parse_plain_text(text: str, filename: str) -> List[Dict[str, str]]:
    text = text.strip()
    if not text:
        return []
    return [{"title": Path(filename or "").stem or "Imported note", "content": text}]


def _parse_markdown_documents(text: str) -> List[Dict[str, str]]:
    lines = text.splitlines()
    docs: List[Dict[str, str]] = []
    buffer: List[str] = []
    current_title: Optional[str] = None

    for line in lines:
        heading = re.match(r"^(#{1,6})\s+(.*)$", line.strip())
        if heading:
            if buffer:
                docs.append({"title": current_title or "Section", "content": "\n".join(buffer).strip()})
                buffer = []
            current_title = heading.group(2).strip()
        else:
            buffer.append(line)

    if buffer:
        docs.append({"title": current_title or "Section", "content": "\n".join(buffer).strip()})

    docs = [d for d in docs if d["content"]]
    if not docs and text.strip():
        docs.append({"title": "Imported note", "content": text.strip()})
    return docs


def _parse_csv_documents(data: bytes, filename: str) -> List[Dict[str, str]]:
    df = pd.read_csv(io.BytesIO(data)).fillna("")
    if df.empty:
        return []
    docs: List[Dict[str, str]] = []
    for idx, row in df.iterrows():
        title = str(row.get("title") or row.get("name") or f"Row {idx + 1}").strip()
        content = str(row.get("content") or row.get("text") or "").strip()
        if not content:
            extra_parts = []
            for col in df.columns:
                if col in {"title", "name", "content", "text"}:
                    continue
                value = str(row.get(col) or "").strip()
                if value:
                    extra_parts.append(f"{col}: {value}")
            content = "\n".join(extra_parts)
        if content:
            docs.append({"title": title or f"Row {idx + 1}", "content": content})
    return docs


def _parse_docx_documents(data: bytes, filename: str) -> List[Dict[str, str]]:
    document = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    if not text:
        return []
    title = document.core_properties.title or Path(filename or "").stem or "DOCX document"
    return [{"title": title, "content": text}]


def _parse_pptx_documents(data: bytes, filename: str) -> List[Dict[str, str]]:
    presentation = Presentation(io.BytesIO(data))
    docs: List[Dict[str, str]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        content = "\n".join(t for t in texts if t)
        if content:
            title = texts[0] if texts else f"Slide {idx}"
            docs.append({"title": title, "content": content})
    if not docs:
        aggregated: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    aggregated.append(shape.text.strip())
        if aggregated:
            title = Path(filename or "").stem or "PPTX document"
            docs.append({"title": title, "content": "\n".join(aggregated)})
    return docs


def _parse_pdf_document(data: bytes, filename: str) -> List[Dict[str, str]]:
    raw_text = extract_pdf_text(io.BytesIO(data))
    pages_raw = [page.strip() for page in raw_text.split("\f") if page.strip()]
    if not pages_raw:
        return []

    page_lines: List[List[str]] = []
    for page in pages_raw:
        lines = [line.strip() for line in page.splitlines() if line.strip()]
        if lines:
            page_lines.append(lines)

    if not page_lines:
        return []

    # detect repeating headers/footers (first/last line that appear on majority pages)
    header_counter = Counter(lines[0] for lines in page_lines if lines)
    footer_counter = Counter(lines[-1] for lines in page_lines if lines)
    threshold = max(2, len(page_lines) // 2)
    header_texts = {text for text, count in header_counter.items() if count >= threshold}
    footer_texts = {text for text, count in footer_counter.items() if count >= threshold}

    docs: List[Dict[str, str]] = []
    base_title = Path(filename or "").stem or "PDF document"
    for idx, lines in enumerate(page_lines, start=1):
        filtered: List[str] = []
        for i, line in enumerate(lines):
            if i == 0 and line in header_texts:
                continue
            if i == len(lines) - 1 and line in footer_texts:
                continue
            filtered.append(line)
        chunks = _split_paragraphs("\n".join(filtered).strip())
        for chunk_idx, chunk in enumerate(chunks, start=1):
            docs.append(
                {
                    "title": f"{base_title} - Page {idx} - Part {chunk_idx}",
                    "content": chunk,
                }
            )

    if not docs:
        docs.append({"title": base_title, "content": "\n\n".join(pages_raw)})
    return docs


def _split_paragraphs(text: str, max_chars: int = 1200) -> List[str]:
    """
    Split text by blank line / sentence boundaries while enforcing max length.
    """
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: List[str] = []

    for para in paragraphs:
        if len(para) <= max_chars:
            chunks.append(para)
            continue
        sentences = re.split(r"(?<=[。．.!?])\s+", para)
        current = ""
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            if len(current) + len(sentence) + 1 <= max_chars:
                current = f"{current} {sentence}".strip()
            else:
                if current:
                    chunks.append(current)
                current = sentence
        if current:
            chunks.append(current)

    return chunks or [text]

